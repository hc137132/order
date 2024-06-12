import pill
pil=pill.Pill()
import random
from PIL import Image
# print(pil.xy())
import time
textstr='可以未引入的包自定义一个简单的名字比如上边的比较复杂每次使用都要用这么长的名字' \
                '所以可以使用定义一个简单的名称但是使用后前边的模块名就不再起作用了'
print(time.time())
#不重复随机抽取4个str
# index=random.sample(range(len(textstr)-1), k=4)
# text=[textstr[index[0]],textstr[index[1]],textstr[index[2]],textstr[index[3]]]
# print(pil.xy())
# os.listdir()用于返回指定的文件夹下包含的文件或文件夹名字的列表，这个列表按字母顺序排序。

import os

# filePath = 'D:\\aaaaa'
# print(os.listdir(filePath))
# for a in os.listdir(filePath):
#     im = Image.open(filePath+'\\'+a)
#     # 看看图片的宽和高
#     w, h = im.size
#     # print(w, h)
#     # 可以用画图工具看需要裁剪的位置
#     # crop（x1，y1，x2，y2） 裁剪的是矩形  左上（x1，y1）  到右下（x2，y2）
#     cropim = im.crop((w-280, h-180, w, h))
#     # 保存，也是保存在当前目录
#     cropim.save(filePath+'\\'+a)

# data={'aa':1}
# data1={'aa':2}
# set=SetToken()
# print(set.Token(data,data1))
import jwt
import datetime
print(datetime.datetime.utcnow())
# 密钥，用于签名和验证JWT
secret_key = 'your-secret-key'

# 载荷，即JWT中包含的信息
payload = {
    'user_id': 123,
    'username': 'john_doe',
    'exp': int(time.time())+3600*2,

}

# 生成JWT
token = jwt.encode(payload, secret_key, algorithm='HS256')
print('Generated JWT:', token)

import jwt
import datetime

# 密钥，用于验证JWT
secret_key = 'hecheng'

# 待验证的JWT
received_token = token

# try:
#     # 验证JWT
#     decoded_payload = jwt.decode(received_token, secret_key, algorithms=['HS256'])
#     print('Decoded Payload:', decoded_payload)
# except jwt.ExpiredSignatureError:
#     print('JWT has expired.')#过期
# except jwt.InvalidTokenError:
#     print('Invalid JWT.')

settoken=pill.SetToken()
# print(settoken.Encrypt(payload))
# token=settoken.Encrypt(payload)
# print(settoken.Decrypt(token))
# print(settoken.Veri_Token(token))
token='eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJ1c2VyIjoiMTEyQHFxLmNvbSsImV4cCI6MTcwOTY1NTg3OX0.WT4cZLwrJ-odKyUTFlTC3xzsqgNmqUapa8ToRVHhB8s '

token='eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJ1c2VyIjoiMTEyQHFxLmNvbSIsImV4cCI6MTcwOTY1NTg3OX0.WT4cZLwrJ-odKyUTFlTC3xzsqgNmqUapa8ToRVHhB8s'

# print(settoken.Decrypt(token))
# print(os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))),'static','avatars'))
# token.split('')
# for key in ['avatar', 'introduction', 'nickname', 'phonenumber']:
#
#     match key:
#         case 'avatar':
#             print(1)
#         case 'introduction':
#             print(2)
#         case 'nickname':
#
#             print(3)
#         case _:
#             print(key)
#             pass
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from io import BytesIO
from bs4 import BeautifulSoup



uploaded_file = r'D:\hc312612\图片\Screenshots\python函数\河楼--2023年度工作总结.docx'

file_extension= uploaded_file.split('.')[-1].lower()
html_content=''
if file_extension == 'docx':
    doc = Document(uploaded_file)
    html_content = '<html><head></head><body>'
    for para in doc.paragraphs:
        html_content += f'<p>{para.text}</p>'
    html_content += '</body></html>'
elif file_extension == 'pptx':
    prs = Presentation(uploaded_file)
    html_content = '<html><head></head><body>'
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                html_content += f'<p>{shape.text}</p>'
    html_content += '</body></html>'
elif file_extension == 'xlsx':
    wb = load_workbook(uploaded_file)
    html_content = '<html><head></head><body>'
    for sheet in wb.sheetnames:
        html_content += f'<h2>{sheet}</h2>'
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            html_content += '<p>'
            for cell in row:
                html_content += f'{cell} '
            html_content += '</p>'
    html_content += '</body></html>'
else:
    pass


print(html_content)